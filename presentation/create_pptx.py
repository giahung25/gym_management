"""Script tạo PowerPoint báo cáo Gym Management System"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# ── Màu sắc theo theme.py ─────────────────────────────────────────────────────
ORANGE     = RGBColor(0xF9, 0x73, 0x16)  # #F97316
ORANGE_L   = RGBColor(0xFE, 0xD7, 0xAA)  # #FED7AA
SIDEBAR    = RGBColor(0x1C, 0x1C, 0x2E)  # #1C1C2E
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x6B, 0x72, 0x80)
GRAY_L     = RGBColor(0xE5, 0xE7, 0xEB)
TEXT_PRI   = RGBColor(0x11, 0x18, 0x27)
GREEN      = RGBColor(0x22, 0xC5, 0x5E)
RED        = RGBColor(0xEF, 0x44, 0x44)
BLUE       = RGBColor(0x3B, 0x82, 0xF6)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
BG         = RGBColor(0xF5, 0xF5, 0xF5)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # Blank layout

# ── Helper functions ───────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, color, transparency=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=TEXT_PRI,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=16, color=TEXT_PRI, bullet_color=ORANGE, spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after  = Pt(4)
        run = p.add_run()
        run.text = f"▶  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


def slide_header(slide, title, subtitle=None,
                 bg_color=SIDEBAR, accent=ORANGE):
    """Dải header trên cùng mỗi slide (trừ slide 1)"""
    # Dải nền header
    add_rect(slide, 0, 0, prs.slide_width, Inches(1.2), bg_color)
    # Thanh accent bên trái
    add_rect(slide, 0, 0, Inches(0.08), Inches(1.2), accent)
    # Title
    add_text_box(slide, title,
                 Inches(0.3), Inches(0.12),
                 Inches(10), Inches(0.7),
                 font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.3), Inches(0.78),
                     Inches(10), Inches(0.36),
                     font_size=14, color=ORANGE_L)


def slide_footer(slide, page_num, total=10):
    """Footer mỗi slide"""
    add_rect(slide, 0, Inches(7.1), prs.slide_width, Inches(0.4), SIDEBAR)
    add_text_box(slide, "GymAdmin Management System  •  2026",
                 Inches(0.3), Inches(7.1), Inches(8), Inches(0.4),
                 font_size=10, color=GRAY, italic=True)
    add_text_box(slide, f"{page_num} / {total}",
                 Inches(12.2), Inches(7.1), Inches(1), Inches(0.4),
                 font_size=10, color=ORANGE, bold=True, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Trang bìa
# ═══════════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(blank_layout)

# Background gradient simulation (2 rectangles)
add_rect(slide1, 0, 0, prs.slide_width, prs.slide_height, SIDEBAR)
add_rect(slide1, Inches(8.5), 0, Inches(4.83), prs.slide_height,
         RGBColor(0x2D, 0x2D, 0x44))

# Accent bar trái
add_rect(slide1, 0, 0, Inches(0.12), prs.slide_height, ORANGE)

# Logo "G" hộp vuông
add_rect(slide1, Inches(1.2), Inches(1.4), Inches(1.2), Inches(1.2), ORANGE)
add_text_box(slide1, "G", Inches(1.2), Inches(1.45), Inches(1.2), Inches(1.1),
             font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Title chính
add_text_box(slide1, "GymAdmin",
             Inches(2.7), Inches(1.3), Inches(8), Inches(0.9),
             font_size=48, bold=True, color=WHITE)
add_text_box(slide1, "MANAGEMENT SYSTEM",
             Inches(2.7), Inches(2.15), Inches(8), Inches(0.55),
             font_size=20, bold=True, color=ORANGE)

# Divider
add_rect(slide1, Inches(1.2), Inches(2.95), Inches(9), Inches(0.04), ORANGE)

# Subtitle
add_text_box(slide1, "Báo cáo Kỹ thuật — Hệ thống Quản lý Phòng Gym Desktop",
             Inches(1.2), Inches(3.15), Inches(10), Inches(0.6),
             font_size=20, color=GRAY_L, italic=True)

# Info block
info_lines = [
    "🛠  Stack: Python + Flet + SQLite3",
    "📅  Ngày báo cáo: 2026-03-20",
    "📊  Trạng thái: MVP 80% hoàn thiện",
]
add_bullet_box(slide1, info_lines,
               Inches(1.2), Inches(3.9), Inches(10), Inches(1.8),
               font_size=16, color=GRAY_L, bullet_color=ORANGE)

# Version tag
add_rect(slide1, Inches(1.2), Inches(5.8), Inches(2.5), Inches(0.5), ORANGE)
add_text_box(slide1, "Version 1.0.0  •  2026",
             Inches(1.2), Inches(5.82), Inches(2.5), Inches(0.46),
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Tổng quan dự án
# ═══════════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(blank_layout)
add_rect(slide2, 0, 0, prs.slide_width, prs.slide_height, BG)
slide_header(slide2, "Tổng quan dự án", "Gym Management System — Mục tiêu & Stack công nghệ")
slide_footer(slide2, 2)

# Cột trái — Mục tiêu
add_rect(slide2, Inches(0.4), Inches(1.4), Inches(5.8), Inches(5.4), WHITE)
add_rect(slide2, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.5), ORANGE)
add_text_box(slide2, "Mục tiêu hệ thống",
             Inches(0.5), Inches(1.42), Inches(5.6), Inches(0.46),
             font_size=16, bold=True, color=WHITE)

goals = [
    "Quản lý hội viên (CRUD + tìm kiếm + lọc nâng cao)",
    "Quản lý gói tập & đăng ký (subscribe/cancel)",
    "Quản lý thiết bị (CRUD + filter trạng thái)",
    "Dashboard KPI thời gian thực từ DB",
    "Biểu đồ doanh thu 6 tháng",
    "Cảnh báo gói tập sắp hết hạn (7 ngày)",
    "Báo cáo thống kê tổng hợp",
]
add_bullet_box(slide2, goals,
               Inches(0.6), Inches(2.05), Inches(5.4), Inches(4.5),
               font_size=13.5, color=TEXT_PRI)

# Cột phải — Stack công nghệ
add_rect(slide2, Inches(6.7), Inches(1.4), Inches(6.2), Inches(5.4), WHITE)
add_rect(slide2, Inches(6.7), Inches(1.4), Inches(6.2), Inches(0.5), SIDEBAR)
add_text_box(slide2, "Stack công nghệ",
             Inches(6.8), Inches(1.42), Inches(6), Inches(0.46),
             font_size=16, bold=True, color=WHITE)

tech_items = [
    ("Python 3.10+", "Ngôn ngữ chính — thuần, không ORM", GREEN),
    ("Flet 0.82.2", "GUI Framework Desktop (pinned version)", BLUE),
    ("SQLite3", "Database cục bộ — data/gym_db.db", AMBER),
    ("Regex (re)", "Validation số điện thoại & email", GRAY),
    ("uuid4", "Sinh ID duy nhất cho mọi record", ORANGE),
    ("contextlib", "Context Manager an toàn cho DB", SIDEBAR),
]

y = 2.1
for name, desc, color in tech_items:
    add_rect(slide2, Inches(6.8), Inches(y), Inches(0.12), Inches(0.45), color)
    add_text_box(slide2, name,
                 Inches(7.05), Inches(y), Inches(2.2), Inches(0.45),
                 font_size=13, bold=True, color=color)
    add_text_box(slide2, desc,
                 Inches(9.2), Inches(y), Inches(3.5), Inches(0.45),
                 font_size=12, color=GRAY)
    y += 0.7


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Kiến trúc hệ thống
# ═══════════════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(blank_layout)
add_rect(slide3, 0, 0, prs.slide_width, prs.slide_height, BG)
slide_header(slide3, "Kiến trúc hệ thống", "Layered Architecture — 5 tầng phân tách rõ ràng")
slide_footer(slide3, 3)

layers = [
    ("GUI Layer",          "login • dashboard • members • memberships • equipment • reports",  ORANGE, Inches(1.35)),
    ("Services Layer",     "member_svc • membership_svc • equipment_svc  [Validation + Logic]", BLUE,   Inches(2.6)),
    ("Repositories Layer", "member_repo • membership_repo • equipment_repo  [SQL CRUD]",        GREEN,  Inches(3.85)),
    ("Models Layer",       "BaseModel • Member • MembershipPlan • MembershipSubscription • Equipment", AMBER, Inches(5.1)),
    ("Core Layer",         "config.py • database.py • security.py  [SQLite3 + Auth]",          SIDEBAR, Inches(6.35)),
]

for (name, desc, color, top) in layers:
    add_rect(slide3, Inches(0.5), top, Inches(12.33), Inches(1.0), WHITE)
    add_rect(slide3, Inches(0.5), top, Inches(2.5), Inches(1.0), color)
    add_text_box(slide3, name,
                 Inches(0.55), Inches(float(top / 914400) + 0.27),
                 Inches(2.4), Inches(0.5),
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide3, desc,
                 Inches(3.2), Inches(float(top / 914400) + 0.27),
                 Inches(9.4), Inches(0.5),
                 font_size=13, color=TEXT_PRI)
    # Arrow connector (mũi tên xuống giữa các layer)

# Arrows
for y_pos in [2.55, 3.8, 5.05]:
    add_text_box(slide3, "▼",
                 Inches(6.4), Inches(y_pos), Inches(0.5), Inches(0.35),
                 font_size=16, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Database Schema
# ═══════════════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(blank_layout)
add_rect(slide4, 0, 0, prs.slide_width, prs.slide_height, BG)
slide_header(slide4, "Database Schema", "SQLite3 — 4 bảng + 6 indexes + Foreign Keys")
slide_footer(slide4, 4)

tables = [
    ("members", ORANGE, Inches(0.3), [
        "id (TEXT, PK, UUID4)",
        "name (TEXT, NOT NULL)",
        "phone (TEXT, NOT NULL)",
        "email, gender, date_of_birth",
        "address, emergency_contact",
        "photo (TEXT)",
        "created_at, updated_at, is_active",
    ]),
    ("membership_plans", BLUE, Inches(3.6), [
        "id (TEXT, PK, UUID4)",
        "name (TEXT, NOT NULL)",
        "duration_days (INTEGER)",
        "price (REAL, NOT NULL)",
        "description (TEXT)",
        "created_at, updated_at, is_active",
    ]),
    ("subscriptions", GREEN, Inches(6.9), [
        "id (TEXT, PK, UUID4)",
        "member_id (FK → members)",
        "plan_id (FK → membership_plans)",
        "price_paid (REAL)",
        "start_date, end_date (TEXT)",
        "status (active/expired/cancelled)",
        "created_at, updated_at, is_active",
    ]),
    ("equipment", AMBER, Inches(10.2), [
        "id (TEXT, PK, UUID4)",
        "name, category (TEXT)",
        "quantity (INTEGER)",
        "status (working/broken/maintenance)",
        "purchase_date, location (TEXT)",
        "notes (TEXT)",
        "created_at, updated_at, is_active",
    ]),
]

for (tname, color, left, cols) in tables:
    add_rect(slide4, left, Inches(1.35), Inches(3.0), Inches(5.5), WHITE)
    add_rect(slide4, left, Inches(1.35), Inches(3.0), Inches(0.5), color)
    add_text_box(slide4, tname,
                 left + Inches(0.05), Inches(1.37),
                 Inches(2.9), Inches(0.46),
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    y = 1.95
    for col in cols:
        add_text_box(slide4, f"• {col}",
                     left + Inches(0.1), Inches(y),
                     Inches(2.8), Inches(0.55),
                     font_size=11, color=TEXT_PRI)
        y += 0.62

# Relationship labels
add_text_box(slide4, "1 : N",
             Inches(3.2), Inches(3.8), Inches(0.7), Inches(0.4),
             font_size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text_box(slide4, "N : 1",
             Inches(9.9), Inches(3.8), Inches(0.7), Inches(0.4),
             font_size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Indexes note
add_rect(slide4, Inches(0.3), Inches(7.0), Inches(12.73), Inches(0.5), ORANGE_L)
add_text_box(slide4,
             "6 Indexes: idx_members_phone • idx_members_is_active • idx_subs_member_id • idx_subs_plan_id • idx_subs_status • idx_equipment_status",
             Inches(0.4), Inches(7.0), Inches(12.53), Inches(0.5),
             font_size=11, color=SIDEBAR, bold=False, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Models & Data Flow
# ═══════════════════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(blank_layout)
add_rect(slide5, 0, 0, prs.slide_width, prs.slide_height, BG)
slide_header(slide5, "Models & Luồng dữ liệu", "Pure Python objects — Hydration pattern — Soft Delete")
slide_footer(slide5, 5)

# BaseModel box
add_rect(slide5, Inches(4.7), Inches(1.4), Inches(3.9), Inches(1.4), SIDEBAR)
add_text_box(slide5, "BaseModel",
             Inches(4.7), Inches(1.42), Inches(3.9), Inches(0.5),
             font_size=16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_text_box(slide5, "id (UUID4)  •  created_at  •  updated_at  •  is_active\nupdate() • delete() • to_dict()",
             Inches(4.7), Inches(1.92), Inches(3.9), Inches(0.7),
             font_size=12, color=GRAY_L, align=PP_ALIGN.CENTER)

# Child models
children = [
    ("Member", ORANGE, Inches(0.5), "name, phone, email\ngender, dob, photo\naddress, emergency_contact"),
    ("MembershipPlan", BLUE, Inches(4.5), "name\nduration_days\nprice, description"),
    ("MembershipSubscription", GREEN, Inches(7.3), "member_id, plan_id\nprice_paid\nstart_date, end_date\nstatus + cancel() + refresh_status()"),
    ("Equipment", AMBER, Inches(11.0), "name, category\nquantity, status\nlocation, notes\nmark_*()"),
]

for (name, color, left, fields) in children:
    add_rect(slide5, left, Inches(3.5), Inches(2.6), Inches(2.8), WHITE)
    add_rect(slide5, left, Inches(3.5), Inches(2.6), Inches(0.45), color)
    add_text_box(slide5, name,
                 left + Inches(0.05), Inches(3.52),
                 Inches(2.5), Inches(0.41),
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide5, fields,
                 left + Inches(0.08), Inches(4.0),
                 Inches(2.44), Inches(2.0),
                 font_size=11, color=TEXT_PRI)
    # Arrow up to BaseModel
    add_text_box(slide5, "↑ kế thừa",
                 left + Inches(0.4), Inches(3.1),
                 Inches(1.8), Inches(0.35),
                 font_size=10, color=color, italic=True)

# Hydration note
add_rect(slide5, Inches(0.3), Inches(6.6), Inches(12.73), Inches(0.5), BLUE)
add_text_box(slide5,
             "Hydration pattern: _row_to_model() dùng __new__() — không tạo UUID mới, giữ nguyên id/created_at từ DB",
             Inches(0.4), Inches(6.62), Inches(12.53), Inches(0.46),
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — GUI & Navigation
# ═══════════════════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(blank_layout)
add_rect(slide6, 0, 0, prs.slide_width, prs.slide_height, BG)
slide_header(slide6, "GUI & Navigation", "Flet Framework — Router pattern — Component reusability")
slide_footer(slide6, 6)

# Sidebar mock
add_rect(slide6, Inches(0.3), Inches(1.4), Inches(2.0), Inches(5.5), SIDEBAR)
add_text_box(slide6, "Sidebar", Inches(0.3), Inches(1.42),
             Inches(2.0), Inches(0.4), font_size=12, bold=True,
             color=ORANGE, align=PP_ALIGN.CENTER)

nav_items = [
    ("Dashboard", ORANGE, True),
    ("Members", GRAY, False),
    ("Gym Packages", GRAY, False),
    ("Equipment", GRAY, False),
    ("Reports", GRAY, False),
]
y = 2.0
for (label, color, active) in nav_items:
    if active:
        add_rect(slide6, Inches(0.42), Inches(y), Inches(1.76), Inches(0.38), ORANGE)
        add_text_box(slide6, f"● {label}", Inches(0.46), Inches(y + 0.04),
                     Inches(1.68), Inches(0.3), font_size=11, bold=True, color=WHITE)
    else:
        add_text_box(slide6, f"○ {label}", Inches(0.46), Inches(y + 0.04),
                     Inches(1.68), Inches(0.3), font_size=11, color=GRAY_L)
    y += 0.5

# Header mock
add_rect(slide6, Inches(2.5), Inches(1.4), Inches(10.5), Inches(0.7), WHITE)
add_rect(slide6, Inches(2.5), Inches(1.9), Inches(10.5), Inches(0.02), GRAY_L)
add_text_box(slide6, "🔍 Search...          🔔  [A] Admin User  ▼",
             Inches(2.6), Inches(1.42), Inches(10.3), Inches(0.46),
             font_size=13, color=GRAY)

# Content area
add_rect(slide6, Inches(2.5), Inches(2.15), Inches(10.5), Inches(4.7), WHITE)
add_text_box(slide6, "Content Area — scroll=AUTO",
             Inches(2.6), Inches(2.2), Inches(10.3), Inches(0.4),
             font_size=12, color=GRAY, italic=True)

# Router pattern
route_items = [
    ("navigate('login')",     "LoginScreen"),
    ("navigate('dashboard')", "DashboardScreen"),
    ("navigate('members')",   "MembersScreen"),
    ("navigate('packages')",  "MembershipsScreen"),
    ("navigate('equipment')", "EquipmentScreen"),
    ("navigate('reports')",   "ReportsScreen"),
]
y = 2.7
for (call, screen) in route_items:
    add_text_box(slide6, call,
                 Inches(2.7), Inches(y), Inches(3.5), Inches(0.38),
                 font_size=11.5, color=BLUE, bold=True)
    add_text_box(slide6, f"→  {screen}",
                 Inches(6.1), Inches(y), Inches(4), Inches(0.38),
                 font_size=11.5, color=GREEN)
    y += 0.5

# Key behaviors note
add_rect(slide6, Inches(2.5), Inches(5.8), Inches(10.5), Inches(1.0), ORANGE_L)
add_text_box(slide6,
             "Mỗi navigate(): page.overlay.clear() → xóa dialog cũ\n"
             "page.on_search_change = None → reset header search callback\n"
             "page.navigate = navigate → monkey-patching, Lazy import tránh circular dependency",
             Inches(2.6), Inches(5.82), Inches(10.3), Inches(0.96),
             font_size=12, color=SIDEBAR)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Dashboard & Tính năng chính
# ═══════════════════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(blank_layout)
add_rect(slide7, 0, 0, prs.slide_width, prs.slide_height, BG)
slide_header(slide7, "Dashboard & Tính năng chính", "KPI Cards — Revenue Chart — Expiring Soon — Active Growth")
slide_footer(slide7, 7)

# KPI Cards mock (4 cards)
kpi_data = [
    ("Tổng hội viên", "248", "+12 tháng này", GREEN),
    ("Sắp hết hạn", "7", "7 ngày tới", AMBER),
    ("Doanh thu tháng", "45,000,000đ", "Năm: 380M", GREEN),
    ("Cần bảo trì", "3", "Thiết bị", RED),
]
x = 0.3
for (label, val, badge, color) in kpi_data:
    add_rect(slide7, Inches(x), Inches(1.4), Inches(3.0), Inches(1.5), WHITE)
    add_rect(slide7, Inches(x), Inches(1.4), Inches(3.0), Inches(0.08), color)
    add_text_box(slide7, val,
                 Inches(x + 0.1), Inches(1.6),
                 Inches(2.8), Inches(0.6),
                 font_size=28, bold=True, color=TEXT_PRI)
    add_text_box(slide7, label,
                 Inches(x + 0.1), Inches(2.15),
                 Inches(2.8), Inches(0.4),
                 font_size=11, color=GRAY)
    add_text_box(slide7, badge,
                 Inches(x + 0.1), Inches(2.5),
                 Inches(2.8), Inches(0.3),
                 font_size=10, color=color, bold=True)
    x += 3.23

# Features list (left)
add_rect(slide7, Inches(0.3), Inches(3.15), Inches(6.0), Inches(3.8), WHITE)
add_rect(slide7, Inches(0.3), Inches(3.15), Inches(6.0), Inches(0.45), SIDEBAR)
add_text_box(slide7, "Tính năng Dashboard",
             Inches(0.3), Inches(3.17), Inches(6.0), Inches(0.41),
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

features = [
    "✅ KPI Cards với badge màu semantic (xanh/vàng/đỏ)",
    "✅ Revenue Bar Chart — 6 tháng gần nhất (custom Flet)",
    "✅ Active Growth Chart — ProgressBar theo gói tập",
    "✅ Recent Member Activity — 5 hội viên mới nhất",
    "✅ Gym Packages Section — badge POPULAR tự động",
    "✅ Equipment Status — wear level progress bar",
    "✅ Expiring Soon Section — cảnh báo badge đỏ ≤3 ngày",
    "✅ Dữ liệu thật từ DB — không dùng mock data",
]
add_bullet_box(slide7, features,
               Inches(0.4), Inches(3.7), Inches(5.8), Inches(3.0),
               font_size=12, color=TEXT_PRI)

# Revenue chart mini
add_rect(slide7, Inches(6.6), Inches(3.15), Inches(6.4), Inches(3.8), WHITE)
add_rect(slide7, Inches(6.6), Inches(3.15), Inches(6.4), Inches(0.45), ORANGE)
add_text_box(slide7, "Revenue Overview — Last 6 months",
             Inches(6.6), Inches(3.17), Inches(6.4), Inches(0.41),
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Bar chart visualization
bar_heights = [0.4, 0.6, 0.9, 0.5, 0.75, 1.2]
bar_labels  = ["T10", "T11", "T12", "T1", "T2", "T3"]
bar_colors  = [ORANGE_L]*5 + [ORANGE]
bx = 6.9
for i, (bh, bl, bc) in enumerate(zip(bar_heights, bar_labels, bar_colors)):
    bottom_y = 6.45
    add_rect(slide7, Inches(bx), Inches(bottom_y - bh), Inches(0.6), Inches(bh), bc)
    add_text_box(slide7, bl,
                 Inches(bx), Inches(6.5),
                 Inches(0.6), Inches(0.3),
                 font_size=10, color=GRAY, align=PP_ALIGN.CENTER)
    bx += 0.82


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Members & Memberships Screen
# ═══════════════════════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(blank_layout)
add_rect(slide8, 0, 0, prs.slide_width, prs.slide_height, BG)
slide_header(slide8, "Members & Memberships", "CRUD đầy đủ — Filter nâng cao — Lịch sử gói tập — Hủy subscription")
slide_footer(slide8, 8)

# Members — trái
add_rect(slide8, Inches(0.3), Inches(1.4), Inches(6.1), Inches(5.6), WHITE)
add_rect(slide8, Inches(0.3), Inches(1.4), Inches(6.1), Inches(0.5), ORANGE)
add_text_box(slide8, "Màn hình Hội viên (members.py)",
             Inches(0.35), Inches(1.42), Inches(6.0), Inches(0.46),
             font_size=14, bold=True, color=WHITE)

members_features = [
    "Bảng hội viên: Avatar (initials + màu hash ID) + Tên/SĐT + Email + Giới tính",
    "Tìm kiếm realtime: theo tên, SĐT, email (on_change → member_repo.search())",
    "Filter giới tính: Dropdown với enable_filter + editable + leading_icon",
    "Filter gói tập: Tất cả / Đang active / Không active",
    "Dialog Add/Edit: Reuse 1 dialog — phân biệt bằng selected_member['obj']",
    "Dialog Chi tiết: Lịch sử toàn bộ subscription với badge màu status",
    "Confirm Dialog trước khi xóa — tránh xóa nhầm",
    "Validation: regex SĐT [0-9+\\-\\s]{7,15} + email format",
    "Header search bar kết nối: page.on_search_change = callback",
    "Lambda closure đúng: lambda e, member=m: open_edit(member)",
]
add_bullet_box(slide8, members_features,
               Inches(0.4), Inches(2.0), Inches(5.9), Inches(4.7),
               font_size=11.5, color=TEXT_PRI)

# Memberships — phải
add_rect(slide8, Inches(6.7), Inches(1.4), Inches(6.3), Inches(5.6), WHITE)
add_rect(slide8, Inches(6.7), Inches(1.4), Inches(6.3), Inches(0.5), BLUE)
add_text_box(slide8, "Màn hình Gói tập (memberships.py)",
             Inches(6.75), Inches(1.42), Inches(6.2), Inches(0.46),
             font_size=14, bold=True, color=WHITE)

memberships_features = [
    "2 Tabs: 'Gói tập' (CRUD Plans) + 'Đăng ký' (Subscriptions)",
    "on_tab_change: load đúng data cho tab đang xem",
    "CRUD Plans: tên, số ngày, giá, mô tả — validate đầy đủ",
    "Đăng ký gói tập: chọn member + plan từ Dropdown",
    "Auto-populate giá gốc khi không nhập giá tùy chỉnh",
    "Nút 'Hủy' chỉ hiện cho subscription STATUS_ACTIVE",
    "Confirm Dialog trước khi hủy subscription",
    "Auto-expire: gọi auto_expire_subscriptions() khi refresh",
    "Status badges: active(xanh) / expired(vàng) / cancelled(đỏ)",
    "Ternary spacer: ft.Container(width=46) giữ layout đều",
]
add_bullet_box(slide8, memberships_features,
               Inches(6.8), Inches(2.0), Inches(6.1), Inches(4.7),
               font_size=11.5, color=TEXT_PRI)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Equipment & Reports
# ═══════════════════════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(blank_layout)
add_rect(slide9, 0, 0, prs.slide_width, prs.slide_height, BG)
slide_header(slide9, "Equipment & Reports", "Filter trạng thái — Summary stats — KPI báo cáo — Rebuild pattern")
slide_footer(slide9, 9)

# Equipment — trái
add_rect(slide9, Inches(0.3), Inches(1.4), Inches(6.1), Inches(5.6), WHITE)
add_rect(slide9, Inches(0.3), Inches(1.4), Inches(6.1), Inches(0.5), AMBER)
add_text_box(slide9, "Màn hình Thiết bị (equipment.py)",
             Inches(0.35), Inches(1.42), Inches(6.0), Inches(0.46),
             font_size=14, bold=True, color=WHITE)

eq_features = [
    "Bảng thiết bị: Tên + Loại + SL + Trạng thái + Ngày mua + Actions",
    "Filter Buttons: Tất cả / Hoạt động / Bảo trì / Hỏng",
    "ElevatedButton (cam) = filter đang chọn — OutlinedButton = còn lại",
    "STATUS_COLORS dict: working→xanh / broken→đỏ / maintenance→vàng",
    "Summary text: Tổng / Hoạt động / Bảo trì / Hỏng — cập nhật tự động",
    "Dialog: 7 fields — tên, loại, SL, trạng thái, ngày mua, vị trí, ghi chú",
    "update_equipment(**kwargs): setattr dynamic — không cần liệt kê field",
    "Soft delete: is_active = 0 thay vì DELETE FROM equipment",
    "State machine: STATUS_WORKING ↔ MAINTENANCE ↔ BROKEN",
]
add_bullet_box(slide9, eq_features,
               Inches(0.4), Inches(2.0), Inches(5.9), Inches(4.8),
               font_size=12, color=TEXT_PRI)

# Reports — phải
add_rect(slide9, Inches(6.7), Inches(1.4), Inches(6.3), Inches(5.6), WHITE)
add_rect(slide9, Inches(6.7), Inches(1.4), Inches(6.3), Inches(0.5), GREEN)
add_text_box(slide9, "Màn hình Báo cáo (reports.py)",
             Inches(6.75), Inches(1.42), Inches(6.2), Inches(0.46),
             font_size=14, bold=True, color=WHITE)

rpt_features = [
    "5 KPI Cards: Tổng / Active / Mới hội viên + Doanh thu tháng/năm",
    "Equipment section: 3 cột — Hoạt động / Bảo trì / Hỏng",
    "Expiring section: Danh sách gói sắp hết hạn trong 7 ngày",
    "Countdown badge: đỏ nếu ≤3 ngày, vàng nếu >3 ngày",
    "build_content() + refresh() pattern: rebuild toàn bộ UI khi refresh",
    "Nút 'Làm mới': reload data từ DB + render lại",
    "Empty state: Hiện text nếu không có gói sắp hết hạn",
    "get_revenue_stats(): monthly / yearly / total revenue",
    "Dữ liệu từ services — không query DB trực tiếp",
]
add_bullet_box(slide9, rpt_features,
               Inches(6.8), Inches(2.0), Inches(6.1), Inches(4.8),
               font_size=12, color=TEXT_PRI)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Đánh giá & Kế hoạch
# ═══════════════════════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(blank_layout)
add_rect(slide10, 0, 0, prs.slide_width, prs.slide_height, BG)
slide_header(slide10, "Đánh giá & Kế hoạch tiếp theo", "Điểm mạnh — Cần cải thiện — Roadmap")
slide_footer(slide10, 10)

# Strengths
add_rect(slide10, Inches(0.3), Inches(1.4), Inches(4.0), Inches(5.6), WHITE)
add_rect(slide10, Inches(0.3), Inches(1.4), Inches(4.0), Inches(0.5), GREEN)
add_text_box(slide10, "✅ Điểm mạnh",
             Inches(0.35), Inches(1.42), Inches(3.9), Inches(0.46),
             font_size=15, bold=True, color=WHITE)

strengths = [
    "Kiến trúc 5 tầng rõ ràng",
    "Design System nhất quán (theme.py)",
    "SQL Injection safe (parameterized ?)",
    "Transaction safe (context manager)",
    "Soft Delete — không mất dữ liệu",
    "Validation tập trung ở Service",
    "Observer pattern (header search)",
    "Closure đúng trong lambda/loop",
    "Foreign Key + 6 Indexes",
    "Kết nối DB thật — không mock",
]
add_bullet_box(slide10, strengths,
               Inches(0.4), Inches(2.0), Inches(3.8), Inches(4.8),
               font_size=12, color=GREEN)

# Issues
add_rect(slide10, Inches(4.6), Inches(1.4), Inches(4.0), Inches(5.6), WHITE)
add_rect(slide10, Inches(4.6), Inches(1.4), Inches(4.0), Inches(0.5), AMBER)
add_text_box(slide10, "⚠️ Cần cải thiện",
             Inches(4.65), Inches(1.42), Inches(3.9), Inches(0.46),
             font_size=15, bold=True, color=WHITE)

issues = [
    "Plaintext password (cần bcrypt)",
    "Không có unit tests",
    "Thiếu pagination",
    "N+1 query khi filter sub status",
    "Thiếu loading/empty states",
    "Equipment filter visual feedback",
    "Stub files chưa dọn",
    "Header search chỉ hoạt động ở Members",
    "Footer version hardcode",
    "DatePicker thay TextField cho ngày",
]
add_bullet_box(slide10, issues,
               Inches(4.7), Inches(2.0), Inches(3.8), Inches(4.8),
               font_size=12, color=AMBER)

# Roadmap
add_rect(slide10, Inches(8.9), Inches(1.4), Inches(4.1), Inches(5.6), WHITE)
add_rect(slide10, Inches(8.9), Inches(1.4), Inches(4.1), Inches(0.5), BLUE)
add_text_box(slide10, "🗺  Roadmap",
             Inches(8.95), Inches(1.42), Inches(4.0), Inches(0.46),
             font_size=15, bold=True, color=WHITE)

roadmap = [
    "Tuần 4:",
    "  • Biểu đồ doanh thu theo tháng",
    "  • Toast / Snackbar notifications",
    "Tuần 5:",
    "  • Viết unit tests (pytest)",
    "  • Xuất báo cáo CSV",
    "  • Upload ảnh hội viên",
    "Tuần 6+:",
    "  • Pagination / lazy load",
    "  • Password hashing (bcrypt)",
    "  • Dọn stub files",
]
add_bullet_box(slide10, roadmap,
               Inches(9.0), Inches(2.0), Inches(3.9), Inches(4.8),
               font_size=12, color=BLUE)

# Bottom summary bar
add_rect(slide10, Inches(0.3), Inches(7.1), Inches(12.73), Inches(0.38), ORANGE)
add_text_box(slide10,
             "Tổng thể: 23 files  •  ~3,228 dòng code  •  5 màn hình  •  4 bảng DB  •  80% hoàn thiện  →  Sẵn sàng MVP",
             Inches(0.4), Inches(7.12), Inches(12.53), Inches(0.34),
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(__file__), "GymAdmin_Presentation.pptx")
prs.save(output_path)
print(f"[OK] Da tao: {output_path}")
